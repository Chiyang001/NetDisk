import os
import shutil
import time
import shortuuid
import zipfile
import tempfile
from datetime import datetime, timedelta
from functools import wraps
from flask import Flask, render_template, request, send_from_directory, jsonify, abort, send_file, session, redirect, url_for, after_this_request
from flask_sqlalchemy import SQLAlchemy
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
from PIL import Image
import io
import threading

app = Flask(__name__)

# --- é…ç½® ---
BASE_DIR = os.path.abspath(os.path.dirname(__file__))
STORAGE_DIR = os.path.join(BASE_DIR, 'storage')
FOLDERZIP_DIR = os.path.join(BASE_DIR, 'folderzip')
STATIC_DIR = os.path.join(BASE_DIR, 'static')
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///disk.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024 * 1024  # 16GB æœ€å¤§ä¸Šä¼ å¤§å°
app.secret_key = 'your_secret_key_here' # ç”¨äºSessionåŠ å¯†

# Session é…ç½® - é˜²æ­¢ä¸‹è½½æ—¶ session ä¸¢å¤±
app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(days=7)
app.config['SESSION_REFRESH_EACH_REQUEST'] = True

# é»˜è®¤ç™»å½•å¯†ç ï¼ˆé¦–æ¬¡è¿è¡Œæ—¶ä½¿ç”¨ï¼‰
DEFAULT_PASSWORD = '123456'

# ç‰ˆæœ¬ä¿¡æ¯
VERSION = 'Ver.2026-0101Beta'
AUTHOR = 'ç‚½é˜³001'
GITHUB_URL = 'https://github.com/Chiyang001?tab=repositories'
BILIBILI_URL = 'https://space.bilibili.com/404891612'
PROJECT_URL = 'https://github.com/Chiyang001/NetDisk'

if not os.path.exists(STORAGE_DIR):
    os.makedirs(STORAGE_DIR)

if not os.path.exists(FOLDERZIP_DIR):
    os.makedirs(FOLDERZIP_DIR)

if not os.path.exists(STATIC_DIR):
    os.makedirs(STATIC_DIR)

db = SQLAlchemy(app)

# --- æ•°æ®åº“æ¨¡å‹ï¼šåˆ†äº«é“¾æ¥ ---
class ShareLink(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    token = db.Column(db.String(20), unique=True, nullable=False)
    file_path = db.Column(db.Text, nullable=False) # ç›¸å¯¹è·¯å¾„ï¼Œå¤šä¸ªæ–‡ä»¶ç”¨ | åˆ†éš”
    expire_at = db.Column(db.DateTime, nullable=True) # Noneè¡¨ç¤ºæ°¸ä¹…
    created_at = db.Column(db.DateTime, default=datetime.now)
    is_batch = db.Column(db.Boolean, default=False) # æ˜¯å¦ä¸ºæ‰¹é‡åˆ†äº«

# --- æ•°æ®åº“æ¨¡å‹ï¼šç³»ç»Ÿè®¾ç½® ---
class Settings(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    key = db.Column(db.String(50), unique=True, nullable=False)
    value = db.Column(db.Text, nullable=True)
    updated_at = db.Column(db.DateTime, default=datetime.now, onupdate=datetime.now)

# --- æ•°æ®åº“æ¨¡å‹ï¼šå¯†ç é‡ç½®ä»¤ç‰Œ ---
class PasswordResetToken(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    token = db.Column(db.String(32), unique=True, nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.now)
    expire_at = db.Column(db.DateTime, nullable=False)
    used = db.Column(db.Boolean, default=False)

# åˆå§‹åŒ–æ•°æ®åº“
with app.app_context():
    db.create_all()
    
    # æ£€æŸ¥å¹¶æ·»åŠ æ–°å­—æ®µï¼ˆå…¼å®¹æ—§æ•°æ®åº“ï¼‰
    try:
        # å°è¯•æŸ¥è¯¢ is_batch å­—æ®µï¼Œå¦‚æœä¸å­˜åœ¨ä¼šæŠ›å‡ºå¼‚å¸¸
        ShareLink.query.with_entities(ShareLink.is_batch).first()
    except Exception as e:
        # å­—æ®µä¸å­˜åœ¨ï¼Œéœ€è¦æ·»åŠ 
        print("æ£€æµ‹åˆ°æ•°æ®åº“éœ€è¦æ›´æ–°ï¼Œæ­£åœ¨æ·»åŠ æ–°å­—æ®µ...")
        try:
            with db.engine.connect() as conn:
                # æ·»åŠ  is_batch å­—æ®µ
                conn.execute(db.text("ALTER TABLE share_link ADD COLUMN is_batch BOOLEAN DEFAULT 0"))
                # ä¿®æ”¹ file_path å­—æ®µç±»å‹ä¸º TEXT
                conn.execute(db.text("ALTER TABLE share_link MODIFY COLUMN file_path TEXT"))
                conn.commit()
                print("æ•°æ®åº“æ›´æ–°å®Œæˆ")
        except Exception as alter_error:
            print(f"æ•°æ®åº“æ›´æ–°å¤±è´¥ï¼ˆå¯èƒ½å·²ç»æ›´æ–°è¿‡ï¼‰: {alter_error}")
    
    # åˆå§‹åŒ–é»˜è®¤è®¾ç½®
    if not Settings.query.filter_by(key='password_hash').first():
        default_hash = generate_password_hash(DEFAULT_PASSWORD)
        db.session.add(Settings(key='password_hash', value=default_hash))
        print("\n" + "="*60)
        print("ğŸ” é¦–æ¬¡è¿è¡Œæ£€æµ‹åˆ°ï¼")
        print(f"ğŸ“ é»˜è®¤ç™»å½•å¯†ç ï¼š{DEFAULT_PASSWORD}")
        print("âš ï¸  è¯·ç«‹å³ç™»å½•å¹¶ä¿®æ”¹å¯†ç ä»¥ç¡®ä¿ç³»ç»Ÿå®‰å…¨ï¼")
        print("="*60 + "\n")
    
    if not Settings.query.filter_by(key='theme').first():
        db.session.add(Settings(key='theme', value='light'))
    
    if not Settings.query.filter_by(key='background_type').first():
        db.session.add(Settings(key='background_type', value='image'))
    
    if not Settings.query.filter_by(key='background_image').first():
        db.session.add(Settings(key='background_image', value='bg.png'))
    
    if not Settings.query.filter_by(key='background_color').first():
        db.session.add(Settings(key='background_color', value='#667eea'))
    
    if not Settings.query.filter_by(key='security_question').first():
        db.session.add(Settings(key='security_question', value=''))
    
    if not Settings.query.filter_by(key='security_answer').first():
        db.session.add(Settings(key='security_answer', value=''))
    
    db.session.commit()

# --- è¾…åŠ©å‡½æ•°ï¼šè·å–è®¾ç½® ---
def get_setting(key, default=None):
    setting = Settings.query.filter_by(key=key).first()
    return setting.value if setting else default

def set_setting(key, value):
    setting = Settings.query.filter_by(key=key).first()
    if setting:
        setting.value = value
        setting.updated_at = datetime.now()
    else:
        setting = Settings(key=key, value=value)
        db.session.add(setting)
    db.session.commit()

# --- è¾…åŠ©å‡½æ•°ï¼šéªŒè¯å¯†ç  ---
def verify_password(password):
    password_hash = get_setting('password_hash')
    if password_hash:
        return check_password_hash(password_hash, password)
    return False

# --- è¾…åŠ©å‡½æ•°ï¼šæ£€æŸ¥æ˜¯å¦ä¸ºé»˜è®¤å¯†ç  ---
def is_default_password():
    """æ£€æŸ¥å½“å‰å¯†ç æ˜¯å¦ä¸ºé»˜è®¤å¯†ç """
    password_hash = get_setting('password_hash')
    if password_hash:
        return check_password_hash(password_hash, DEFAULT_PASSWORD)
    return True

# --- è¾…åŠ©å‡½æ•°ï¼šæ¸…ç†è¿‡æœŸçš„ ZIP æ–‡ä»¶ ---
def cleanup_old_zips():
    """åˆ é™¤è¶…è¿‡24å°æ—¶çš„ ZIP æ–‡ä»¶"""
    try:
        now = time.time()
        for filename in os.listdir(FOLDERZIP_DIR):
            filepath = os.path.join(FOLDERZIP_DIR, filename)
            if os.path.isfile(filepath) and filename.endswith('.zip'):
                file_age = now - os.path.getmtime(filepath)
                if file_age > 24 * 3600:  # 24å°æ—¶
                    os.remove(filepath)
                    print(f"å·²åˆ é™¤è¿‡æœŸ ZIP æ–‡ä»¶: {filename}")
    except Exception as e:
        print(f"æ¸…ç† ZIP æ–‡ä»¶å¤±è´¥: {e}")

# --- åå°å®šæ—¶æ¸…ç†ä»»åŠ¡ ---
def schedule_cleanup():
    """æ¯å°æ—¶æ‰§è¡Œä¸€æ¬¡æ¸…ç†ä»»åŠ¡"""
    cleanup_old_zips()
    # è®¾ç½®ä¸‹æ¬¡æ‰§è¡Œ
    threading.Timer(3600, schedule_cleanup).start()

# å¯åŠ¨æ¸…ç†ä»»åŠ¡
schedule_cleanup()

# --- ç™»å½•éªŒè¯è£…é¥°å™¨ ---
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('logged_in'):
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

# --- è¾…åŠ©å‡½æ•°ï¼šå®‰å…¨è·¯å¾„æ£€æŸ¥ ---
def get_safe_path(req_path):
    # é˜²æ­¢ ../ æ”»å‡»ï¼Œç¡®ä¿è·¯å¾„åœ¨ storage ç›®å½•ä¸‹
    if not req_path or req_path.strip() == '/':
        return STORAGE_DIR
    
    # ç§»é™¤å¼€å¤´çš„ /
    req_path = req_path.lstrip('/')
    safe_path = os.path.abspath(os.path.join(STORAGE_DIR, req_path))
    if not safe_path.startswith(STORAGE_DIR):
        raise ValueError("éæ³•è·¯å¾„")
    return safe_path

def get_rel_path(full_path):
    return full_path.replace(STORAGE_DIR, '').replace('\\', '/').lstrip('/')

# --- è¾…åŠ©å‡½æ•°ï¼šåˆ¤æ–­æ–‡ä»¶ç±»å‹ ---
def is_image(filename):
    image_exts = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg', '.ico'}
    return os.path.splitext(filename.lower())[1] in image_exts

def is_video(filename):
    video_exts = {'.mp4', '.webm', '.ogg', '.mov', '.avi', '.mkv', '.flv', '.wmv'}
    return os.path.splitext(filename.lower())[1] in video_exts

def is_audio(filename):
    audio_exts = {'.mp3', '.wav', '.ogg', '.m4a', '.aac', '.flac', '.wma', '.ape', '.opus'}
    return os.path.splitext(filename.lower())[1] in audio_exts

def is_archive(filename):
    archive_exts = {'.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.xz', '.tar.gz', '.tar.bz2', '.tar.xz'}
    lower_name = filename.lower()
    return any(lower_name.endswith(ext) for ext in archive_exts)

def is_office_doc(filename):
    office_exts = {'.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'}
    return os.path.splitext(filename.lower())[1] in office_exts

def is_pdf(filename):
    return os.path.splitext(filename.lower())[1] == '.pdf'

def get_file_type(filename):
    if is_image(filename):
        return 'image'
    elif is_video(filename):
        return 'video'
    elif is_audio(filename):
        return 'audio'
    elif is_archive(filename):
        return 'archive'
    elif is_office_doc(filename):
        return 'office'
    elif is_pdf(filename):
        return 'pdf'
    else:
        return 'file'

# --- è·¯ç”±ï¼šç™»å½•é¡µé¢ ---
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        password = request.form.get('password')
        if verify_password(password):
            # æ£€æŸ¥æ˜¯å¦å·²è®¾ç½®å®‰å…¨é—®é¢˜
            security_question = get_setting('security_question', '')
            if not security_question:
                # æœªè®¾ç½®å®‰å…¨é—®é¢˜ï¼Œæ ‡è®°éœ€è¦è®¾ç½®å¹¶è·³è½¬åˆ°è®¾ç½®é¡µé¢
                session['temp_logged_in'] = True  # ä¸´æ—¶ç™»å½•çŠ¶æ€
                session.permanent = True
                return redirect(url_for('setup_security'))
            
            session['logged_in'] = True
            session.permanent = True  # ä½¿ session æŒä¹…åŒ–
            return redirect(url_for('index'))
        else:
            # è·å–ä¸»é¢˜å’ŒèƒŒæ™¯è®¾ç½®ï¼ˆé”™è¯¯æ—¶ä¹Ÿè¦ä¼ é€’ï¼‰
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('login.html', 
                                 error='å¯†ç é”™è¯¯',
                                 theme=theme, 
                                 bg_type=bg_type, 
                                 bg_image=bg_image, 
                                 bg_color=bg_color,
                                 is_default_password=is_default_password(),
                                 reset_success=False)
    
    # è·å–ä¸»é¢˜å’ŒèƒŒæ™¯è®¾ç½®
    theme = get_setting('theme', 'light')
    bg_type = get_setting('background_type', 'image')
    bg_image = get_setting('background_image', 'bg.png')
    bg_color = get_setting('background_color', '#667eea')
    
    # æ£€æŸ¥æ˜¯å¦æœ‰é‡ç½®æˆåŠŸçš„æ ‡è®°
    reset_success = request.args.get('reset') == 'success'
    
    return render_template('login.html', 
                         theme=theme, 
                         bg_type=bg_type, 
                         bg_image=bg_image, 
                         bg_color=bg_color,
                         is_default_password=is_default_password(),
                         reset_success=reset_success)

# --- è·¯ç”±ï¼šç™»å‡º ---
@app.route('/logout')
def logout():
    session.pop('logged_in', None)
    session.pop('temp_logged_in', None)
    return redirect(url_for('login'))

# --- è·¯ç”±ï¼šè®¾ç½®å®‰å…¨é—®é¢˜ï¼ˆé¦–æ¬¡ç™»å½•å¿…é¡»è®¾ç½®ï¼‰---
@app.route('/setup-security', methods=['GET', 'POST'])
def setup_security():
    # æ£€æŸ¥æ˜¯å¦æœ‰ä¸´æ—¶ç™»å½•çŠ¶æ€
    if not session.get('temp_logged_in'):
        return redirect(url_for('login'))
    
    if request.method == 'POST':
        question = request.form.get('question', '').strip()
        custom_question = request.form.get('custom_question', '').strip()
        answer = request.form.get('answer', '').strip()
        
        # å¦‚æœé€‰æ‹©äº†è‡ªå®šä¹‰é—®é¢˜
        if question == 'custom':
            question = custom_question
        
        if not question or not answer:
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('setup_security.html',
                                 error='é—®é¢˜å’Œç­”æ¡ˆä¸èƒ½ä¸ºç©º',
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color)
        
        if len(answer) < 2:
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('setup_security.html',
                                 error='ç­”æ¡ˆè‡³å°‘éœ€è¦2ä¸ªå­—ç¬¦',
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color)
        
        # ä¿å­˜å®‰å…¨é—®é¢˜
        set_setting('security_question', question)
        set_setting('security_answer', answer)
        
        # æ¸…é™¤ä¸´æ—¶ç™»å½•çŠ¶æ€ï¼Œè®¾ç½®æ­£å¼ç™»å½•çŠ¶æ€
        session.pop('temp_logged_in', None)
        session['logged_in'] = True
        session.permanent = True
        
        return redirect(url_for('index'))
    
    # GET è¯·æ±‚
    theme = get_setting('theme', 'light')
    bg_type = get_setting('background_type', 'image')
    bg_image = get_setting('background_image', 'bg.png')
    bg_color = get_setting('background_color', '#667eea')
    
    return render_template('setup_security.html',
                         theme=theme,
                         bg_type=bg_type,
                         bg_image=bg_image,
                         bg_color=bg_color)

# --- è·¯ç”±ï¼šå¿˜è®°å¯†ç é¡µé¢ ---
@app.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    if request.method == 'POST':
        security_answer = request.form.get('security_answer', '').strip()
        
        # è·å–è®¾ç½®çš„å®‰å…¨é—®é¢˜ç­”æ¡ˆ
        saved_answer = get_setting('security_answer', '')
        
        if not saved_answer:
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('forgot_password.html',
                                 error='ç®¡ç†å‘˜æœªè®¾ç½®å®‰å…¨é—®é¢˜ï¼Œè¯·è”ç³»ç®¡ç†å‘˜é‡ç½®å¯†ç ',
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color,
                                 security_question=get_setting('security_question', ''))
        
        # éªŒè¯ç­”æ¡ˆï¼ˆä¸åŒºåˆ†å¤§å°å†™ï¼‰
        if security_answer.lower() == saved_answer.lower():
            # ç”Ÿæˆé‡ç½®ä»¤ç‰Œ
            token = shortuuid.uuid()
            expire_at = datetime.now() + timedelta(minutes=30)  # 30åˆ†é’Ÿæœ‰æ•ˆæœŸ
            
            reset_token = PasswordResetToken(token=token, expire_at=expire_at)
            db.session.add(reset_token)
            db.session.commit()
            
            # é‡å®šå‘åˆ°é‡ç½®å¯†ç é¡µé¢
            return redirect(url_for('reset_password', token=token))
        else:
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('forgot_password.html',
                                 error='å®‰å…¨é—®é¢˜ç­”æ¡ˆé”™è¯¯',
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color,
                                 security_question=get_setting('security_question', ''))
    
    # GET è¯·æ±‚
    theme = get_setting('theme', 'light')
    bg_type = get_setting('background_type', 'image')
    bg_image = get_setting('background_image', 'bg.png')
    bg_color = get_setting('background_color', '#667eea')
    security_question = get_setting('security_question', '')
    
    if not security_question:
        return render_template('forgot_password.html',
                             error='ç®¡ç†å‘˜æœªè®¾ç½®å®‰å…¨é—®é¢˜ï¼Œè¯·è”ç³»ç®¡ç†å‘˜é‡ç½®å¯†ç ',
                             theme=theme,
                             bg_type=bg_type,
                             bg_image=bg_image,
                             bg_color=bg_color,
                             security_question='')
    
    return render_template('forgot_password.html',
                         theme=theme,
                         bg_type=bg_type,
                         bg_image=bg_image,
                         bg_color=bg_color,
                         security_question=security_question)

# --- è·¯ç”±ï¼šé‡ç½®å¯†ç é¡µé¢ ---
@app.route('/reset-password/<token>', methods=['GET', 'POST'])
def reset_password(token):
    # éªŒè¯ä»¤ç‰Œ
    reset_token = PasswordResetToken.query.filter_by(token=token, used=False).first()
    
    if not reset_token:
        return "é‡ç½®é“¾æ¥æ— æ•ˆæˆ–å·²ä½¿ç”¨", 403
    
    if datetime.now() > reset_token.expire_at:
        return "é‡ç½®é“¾æ¥å·²è¿‡æœŸï¼Œè¯·é‡æ–°ç”³è¯·", 403
    
    if request.method == 'POST':
        new_password = request.form.get('new_password', '').strip()
        confirm_password = request.form.get('confirm_password', '').strip()
        
        if not new_password or len(new_password) < 6:
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('reset_password.html',
                                 error='å¯†ç è‡³å°‘éœ€è¦6ä½',
                                 token=token,
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color)
        
        if new_password != confirm_password:
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('reset_password.html',
                                 error='ä¸¤æ¬¡è¾“å…¥çš„å¯†ç ä¸ä¸€è‡´',
                                 token=token,
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color)
        
        # æ›´æ–°å¯†ç 
        new_hash = generate_password_hash(new_password)
        set_setting('password_hash', new_hash)
        
        # æ ‡è®°ä»¤ç‰Œä¸ºå·²ä½¿ç”¨
        reset_token.used = True
        db.session.commit()
        
        # é‡å®šå‘åˆ°ç™»å½•é¡µé¢
        return redirect(url_for('login') + '?reset=success')
    
    # GET è¯·æ±‚
    theme = get_setting('theme', 'light')
    bg_type = get_setting('background_type', 'image')
    bg_image = get_setting('background_image', 'bg.png')
    bg_color = get_setting('background_color', '#667eea')
    
    return render_template('reset_password.html',
                         token=token,
                         theme=theme,
                         bg_type=bg_type,
                         bg_image=bg_image,
                         bg_color=bg_color)

# --- è·¯ç”±ï¼šè®¾ç½®é¡µé¢ ---
@app.route('/settings')
@login_required
def settings_page():
    # è·å–ä¸»é¢˜å’ŒèƒŒæ™¯è®¾ç½®
    theme = get_setting('theme', 'light')
    bg_type = get_setting('background_type', 'image')
    bg_image = get_setting('background_image', 'bg.png')
    bg_color = get_setting('background_color', '#667eea')
    
    return render_template('settings.html',
                         theme=theme,
                         bg_type=bg_type,
                         bg_image=bg_image,
                         bg_color=bg_color,
                         version=VERSION,
                         author=AUTHOR,
                         github_url=GITHUB_URL,
                         bilibili_url=BILIBILI_URL,
                         project_url=PROJECT_URL)

# --- è·¯ç”±ï¼šé¦–é¡µ (æ–‡ä»¶åˆ—è¡¨) ---
@app.route('/')
@login_required
def index():
    # è·å–å½“å‰è¯·æ±‚çš„ç›¸å¯¹è·¯å¾„ï¼Œé»˜è®¤ä¸ºæ ¹ç›®å½•
    req_path = request.args.get('path', '')
    sort_by = request.args.get('sort', 'name')  # name, time, size
    sort_order = request.args.get('order', 'asc')  # asc, desc
    
    try:
        abs_path = get_safe_path(req_path)
    except:
        return "éæ³•è·¯å¾„", 403

    files_list = []
    if os.path.isdir(abs_path):
        for item in os.listdir(abs_path):
            if item.startswith('.'): continue # éšè—æ–‡ä»¶
            full_item_path = os.path.join(abs_path, item)
            is_dir = os.path.isdir(full_item_path)
            size = os.path.getsize(full_item_path) if not is_dir else 0
            size_bytes = size  # ä¿å­˜åŸå§‹å­—èŠ‚æ•°ç”¨äºæ’åº
            mtime_timestamp = os.path.getmtime(full_item_path)  # ä¿å­˜æ—¶é—´æˆ³ç”¨äºæ’åº
            # è½¬æ¢æ—¶é—´
            mtime = time.strftime('%Y-%m-%d %H:%M', time.localtime(mtime_timestamp))
            
            rel_path_item = os.path.join(req_path, item).replace('\\', '/')
            
            files_list.append({
                'name': item,
                'is_dir': is_dir,
                'size': f"{size/1024/1024:.2f} MB" if not is_dir else "-",
                'size_bytes': size_bytes,
                'mtime': mtime,
                'mtime_timestamp': mtime_timestamp,
                'rel_path': rel_path_item,
                'file_type': 'folder' if is_dir else get_file_type(item)
            })
    
    # æ’åºé€»è¾‘
    if sort_by == 'name':
        files_list.sort(key=lambda x: (not x['is_dir'], x['name'].lower()))
    elif sort_by == 'time':
        files_list.sort(key=lambda x: (not x['is_dir'], x['mtime_timestamp']))
    elif sort_by == 'size':
        files_list.sort(key=lambda x: (not x['is_dir'], x['size_bytes']))
    
    # å€’åº
    if sort_order == 'desc':
        # åˆ†ç¦»æ–‡ä»¶å¤¹å’Œæ–‡ä»¶
        folders = [f for f in files_list if f['is_dir']]
        files = [f for f in files_list if not f['is_dir']]
        # åˆ†åˆ«å€’åº
        folders.reverse()
        files.reverse()
        files_list = folders + files
    
    # è·å–ä¸»é¢˜å’ŒèƒŒæ™¯è®¾ç½®
    theme = get_setting('theme', 'light')
    bg_type = get_setting('background_type', 'image')
    bg_image = get_setting('background_image', 'bg.png')
    bg_color = get_setting('background_color', '#667eea')
    
    return render_template('index.html', 
                         files=files_list, 
                         current_path=req_path,
                         sort_by=sort_by,
                         sort_order=sort_order,
                         theme=theme,
                         bg_type=bg_type,
                         bg_image=bg_image,
                         bg_color=bg_color,
                         version=VERSION,
                         author=AUTHOR,
                         github_url=GITHUB_URL,
                         bilibili_url=BILIBILI_URL)

# --- æ¥å£ï¼šæ“ä½œ (é‡å‘½å, åˆ é™¤, æ–°å»ºæ–‡ä»¶å¤¹) ---
@app.route('/api/operate', methods=['POST'])
@login_required
def operate():
    data = request.json
    action = data.get('action')
    path = data.get('path') # ç›¸å¯¹è·¯å¾„
    
    try:
        abs_path = get_safe_path(path)
        
        if action == 'mkdir':
            new_folder = data.get('name')
            # ä¸ä½¿ç”¨ secure_filenameï¼Œå› ä¸ºå®ƒä¼šç§»é™¤ä¸­æ–‡å­—ç¬¦
            # åªç§»é™¤å±é™©å­—ç¬¦
            new_folder = new_folder.replace('..', '').replace('/', '').replace('\\', '')
            if not new_folder or new_folder.strip() == '':
                return jsonify({'status': 'error', 'msg': 'æ–‡ä»¶å¤¹åç§°æ— æ•ˆ'})
            
            new_folder_path = os.path.join(abs_path, new_folder)
            
            # æ£€æŸ¥æ–‡ä»¶å¤¹æ˜¯å¦å·²å­˜åœ¨
            if os.path.exists(new_folder_path):
                return jsonify({'status': 'error', 'msg': f'æ–‡ä»¶å¤¹ "{new_folder}" å·²å­˜åœ¨'})
            
            os.mkdir(new_folder_path)
            
        elif action == 'delete':
            if os.path.isdir(abs_path):
                shutil.rmtree(abs_path)
            else:
                os.remove(abs_path)
                
        elif action == 'rename':
            new_name = data.get('new_name')
            # ä¸ä½¿ç”¨ secure_filenameï¼Œä¿ç•™ä¸­æ–‡
            new_name = new_name.replace('..', '').replace('/', '').replace('\\', '')
            if not new_name or new_name.strip() == '':
                return jsonify({'status': 'error', 'msg': 'åç§°æ— æ•ˆ'})
            
            parent = os.path.dirname(abs_path)
            new_path = os.path.join(parent, new_name)
            
            # æ£€æŸ¥ç›®æ ‡åç§°æ˜¯å¦å·²å­˜åœ¨
            if os.path.exists(new_path):
                return jsonify({'status': 'error', 'msg': f'åç§° "{new_name}" å·²å­˜åœ¨'})
            
            os.rename(abs_path, new_path)
            
        return jsonify({'status': 'success'})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'msg': str(e)})

# --- æ¥å£ï¼šå¤åˆ¶/ç§»åŠ¨/ç²˜è´´ ---
@app.route('/api/paste', methods=['POST'])
@login_required
def paste():
    data = request.json
    src_path = data.get('src')
    dest_path = data.get('dest') # ç›®æ ‡æ–‡ä»¶å¤¹
    action = data.get('action') # copy æˆ– move

    try:
        abs_src = get_safe_path(src_path)
        abs_dest_folder = get_safe_path(dest_path)
        
        # æ£€æŸ¥æºæ–‡ä»¶æ˜¯å¦å­˜åœ¨
        if not os.path.exists(abs_src):
            return jsonify({'status': 'error', 'msg': 'æºæ–‡ä»¶ä¸å­˜åœ¨'})
        
        # æ£€æŸ¥ç›®æ ‡æ–‡ä»¶å¤¹æ˜¯å¦å­˜åœ¨
        if not os.path.exists(abs_dest_folder):
            return jsonify({'status': 'error', 'msg': 'ç›®æ ‡æ–‡ä»¶å¤¹ä¸å­˜åœ¨'})
        
        filename = os.path.basename(abs_src)
        abs_dest_final = os.path.join(abs_dest_folder, filename)
        
        # å¦‚æœç›®æ ‡å·²å­˜åœ¨ï¼Œè‡ªåŠ¨é‡å‘½å
        if os.path.exists(abs_dest_final):
            base_name, ext = os.path.splitext(filename)
            counter = 1
            while os.path.exists(abs_dest_final):
                new_filename = f"{base_name}_å‰¯æœ¬{counter}{ext}"
                abs_dest_final = os.path.join(abs_dest_folder, new_filename)
                counter += 1

        if action == 'copy':
            if os.path.isdir(abs_src):
                shutil.copytree(abs_src, abs_dest_final)
            else:
                shutil.copy2(abs_src, abs_dest_final)
        elif action == 'move':
            shutil.move(abs_src, abs_dest_final)
            
        return jsonify({'status': 'success'})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'msg': f'æ“ä½œå¤±è´¥: {str(e)}'})

# --- æ¥å£ï¼šä¸Šä¼ æ–‡ä»¶ ---
@app.route('/upload', methods=['POST'])
@login_required
def upload():
    current_path = request.form.get('path', '')
    file = request.files.get('file')
    relative_path = request.form.get('relativePath', '')
    
    try:
        save_dir = get_safe_path(current_path)
        
        if file:
            # å®‰å…¨çš„æ–‡ä»¶åå¤„ç†å‡½æ•°ï¼ˆä¿ç•™ä¸­æ–‡ï¼‰
            def safe_filename(filename):
                # åªç§»é™¤å±é™©å­—ç¬¦ï¼Œä¿ç•™ä¸­æ–‡å’Œå…¶ä»–å­—ç¬¦
                dangerous_chars = ['..', '/', '\\', '\0', '<', '>', ':', '"', '|', '?', '*']
                for char in dangerous_chars:
                    filename = filename.replace(char, '_')
                return filename.strip()
            
            # å¦‚æœæœ‰ç›¸å¯¹è·¯å¾„ï¼ˆæ–‡ä»¶å¤¹ä¸Šä¼ ï¼‰ï¼Œä¿æŒç›®å½•ç»“æ„
            if relative_path and '/' in relative_path:
                # æå–ç›®å½•éƒ¨åˆ†
                path_parts = relative_path.split('/')
                if len(path_parts) > 1:
                    # åˆ›å»ºå­ç›®å½•ï¼ˆä¿ç•™ä¸­æ–‡ç›®å½•åï¼‰
                    safe_parts = [safe_filename(part) for part in path_parts[:-1]]
                    sub_dir = os.path.join(save_dir, *safe_parts)
                    os.makedirs(sub_dir, exist_ok=True)
                    filename = safe_filename(path_parts[-1])
                    file.save(os.path.join(sub_dir, filename))
                else:
                    filename = safe_filename(relative_path)
                    file.save(os.path.join(save_dir, filename))
            else:
                # æ™®é€šæ–‡ä»¶ä¸Šä¼ 
                filename = safe_filename(file.filename)
                file.save(os.path.join(save_dir, filename))
                
        return jsonify({'status': 'success'})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'msg': str(e)})

# --- æ¥å£ï¼šåˆ›å»ºåˆ†äº« ---
@app.route('/api/share', methods=['POST'])
@login_required
def create_share():
    data = request.json
    path = data.get('path')
    paths = data.get('paths')  # æ‰¹é‡åˆ†äº«
    minutes = int(data.get('minutes', 0)) # 0ä»£è¡¨æ°¸ä¹…
    
    token = shortuuid.uuid()[:8] # ç”Ÿæˆ8ä½çŸ­é“¾æ¥
    expire_at = datetime.now() + timedelta(minutes=minutes) if minutes > 0 else None
    
    if paths and len(paths) > 0:
        # æ‰¹é‡åˆ†äº«ï¼šå¤šä¸ªè·¯å¾„ç”¨ | åˆ†éš”
        file_path = '|'.join(paths)
        new_share = ShareLink(token=token, file_path=file_path, expire_at=expire_at, is_batch=True)
    else:
        # å•ä¸ªåˆ†äº«
        new_share = ShareLink(token=token, file_path=path, expire_at=expire_at, is_batch=False)
    
    db.session.add(new_share)
    db.session.commit()
    
    share_url = request.host_url + 's/' + token
    return jsonify({'status': 'success', 'url': share_url})

# --- è¾…åŠ©å‡½æ•°ï¼šæ‰“åŒ…æ–‡ä»¶å¤¹ä¸º ZIP ---
def zip_folder(folder_path, zip_name):
    """å°†æ–‡ä»¶å¤¹æ‰“åŒ…ä¸º ZIP æ–‡ä»¶å¹¶è¿”å›æ–‡ä»¶è·¯å¾„"""
    # ä½¿ç”¨æ—¶é—´æˆ³å’Œéšæœºå­—ç¬¦ä¸²é¿å…æ–‡ä»¶åå†²çª
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    random_str = shortuuid.uuid()[:6]
    safe_zip_name = f"{zip_name}_{timestamp}_{random_str}.zip"
    zip_path = os.path.join(FOLDERZIP_DIR, safe_zip_name)
    
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root, dirs, files in os.walk(folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, folder_path)
                try:
                    zipf.write(file_path, arcname)
                except Exception as e:
                    print(f"æ‰“åŒ…æ–‡ä»¶å¤±è´¥ {file_path}: {e}")
    
    return zip_path

# --- è·¯ç”±ï¼šè®¿é—®åˆ†äº«é“¾æ¥ï¼ˆæ˜¾ç¤ºè¯¦æƒ…é¡µï¼‰---
@app.route('/s/<token>')
def access_share(token):
    link = ShareLink.query.filter_by(token=token).first()
    
    if not link:
        return "é“¾æ¥ä¸å­˜åœ¨æˆ–å·²å¤±æ•ˆ", 404
        
    if link.expire_at and datetime.now() > link.expire_at:
        return "é“¾æ¥å·²è¿‡æœŸ", 403
    
    try:
        # æ£€æŸ¥æ˜¯å¦ä¸ºæ‰¹é‡åˆ†äº«
        if link.is_batch:
            # æ‰¹é‡åˆ†äº« - æ˜¾ç¤ºæ–‡ä»¶åˆ—è¡¨é¡µé¢
            file_paths = link.file_path.split('|')
            file_count = len(file_paths)
            
            # è·å–æ‰€æœ‰æ–‡ä»¶ä¿¡æ¯
            files = []
            total_size = 0
            
            for path in file_paths:
                try:
                    abs_path = get_safe_path(path)
                    if not os.path.exists(abs_path):
                        continue
                    
                    file_name = os.path.basename(abs_path)
                    is_dir = os.path.isdir(abs_path)
                    
                    if is_dir:
                        # è®¡ç®—æ–‡ä»¶å¤¹å¤§å°
                        dir_size = 0
                        for dirpath, dirnames, filenames in os.walk(abs_path):
                            for f in filenames:
                                fp = os.path.join(dirpath, f)
                                if os.path.exists(fp):
                                    dir_size += os.path.getsize(fp)
                        size_bytes = dir_size
                        file_type = 'folder'
                        type_text = 'æ–‡ä»¶å¤¹'
                    else:
                        size_bytes = os.path.getsize(abs_path)
                        file_type = get_file_type(file_name)
                        type_text = {
                            'image': 'å›¾ç‰‡',
                            'video': 'è§†é¢‘',
                            'audio': 'éŸ³é¢‘',
                            'archive': 'å‹ç¼©åŒ…',
                            'pdf': 'PDFæ–‡æ¡£',
                            'office': 'Officeæ–‡æ¡£',
                            'file': 'æ–‡ä»¶'
                        }.get(file_type, 'æ–‡ä»¶')
                    
                    # æ ¼å¼åŒ–å¤§å°
                    if size_bytes < 1024:
                        size_str = f"{size_bytes} B"
                    elif size_bytes < 1024 * 1024:
                        size_str = f"{size_bytes/1024:.2f} KB"
                    else:
                        size_str = f"{size_bytes/1024/1024:.2f} MB"
                    
                    total_size += size_bytes
                    
                    files.append({
                        'name': file_name,
                        'path': path,
                        'size': size_str,
                        'is_dir': is_dir,
                        'file_type': file_type,
                        'type_text': type_text
                    })
                except Exception as e:
                    print(f"å¤„ç†æ–‡ä»¶å¤±è´¥ {path}: {e}")
                    continue
            
            # æ ¼å¼åŒ–æ€»å¤§å°
            if total_size < 1024:
                total_size_str = f"{total_size} B"
            elif total_size < 1024 * 1024:
                total_size_str = f"{total_size/1024:.2f} KB"
            else:
                total_size_str = f"{total_size/1024/1024:.2f} MB"
            
            # æ ¼å¼åŒ–æ—¶é—´
            created_at = link.created_at.strftime('%Y-%m-%d %H:%M')
            expire_at = link.expire_at.strftime('%Y-%m-%d %H:%M') if link.expire_at else None
            
            # è®¡ç®—å‰©ä½™æ—¶é—´
            time_remaining = None
            if link.expire_at:
                remaining = link.expire_at - datetime.now()
                if remaining.days > 0:
                    time_remaining = f"{remaining.days} å¤©"
                elif remaining.seconds > 3600:
                    time_remaining = f"{remaining.seconds // 3600} å°æ—¶"
                elif remaining.seconds > 60:
                    time_remaining = f"{remaining.seconds // 60} åˆ†é’Ÿ"
                else:
                    time_remaining = f"{remaining.seconds} ç§’"
            
            # è·å–ä¸»é¢˜å’ŒèƒŒæ™¯è®¾ç½®
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('batch_share.html',
                                 token=token,
                                 files=files,
                                 file_count=file_count,
                                 total_size=total_size_str,
                                 created_at=created_at,
                                 expire_at=expire_at,
                                 time_remaining=time_remaining,
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color)
        
        else:
            # å•ä¸ªåˆ†äº« - æ˜¾ç¤ºåŸæ¥çš„è¯¦æƒ…é¡µ
            abs_path = get_safe_path(link.file_path)
            
            if not os.path.exists(abs_path):
                return "æ–‡ä»¶æºå·²è¢«åˆ é™¤", 404
            
            # è·å–æ–‡ä»¶ä¿¡æ¯
            file_name = os.path.basename(abs_path)
            is_dir = os.path.isdir(abs_path)
            
            if is_dir:
                file_type = 'folder'
                # è®¡ç®—æ–‡ä»¶å¤¹å¤§å°
                total_size = 0
                for dirpath, dirnames, filenames in os.walk(abs_path):
                    for f in filenames:
                        fp = os.path.join(dirpath, f)
                        if os.path.exists(fp):
                            total_size += os.path.getsize(fp)
                file_size = f"{total_size/1024/1024:.2f} MB"
            else:
                file_type = get_file_type(file_name)
                size_bytes = os.path.getsize(abs_path)
                if size_bytes < 1024:
                    file_size = f"{size_bytes} B"
                elif size_bytes < 1024 * 1024:
                    file_size = f"{size_bytes/1024:.2f} KB"
                else:
                    file_size = f"{size_bytes/1024/1024:.2f} MB"
            
            # æ ¼å¼åŒ–æ—¶é—´
            created_at = link.created_at.strftime('%Y-%m-%d %H:%M')
            expire_at = link.expire_at.strftime('%Y-%m-%d %H:%M') if link.expire_at else None
            
            # è®¡ç®—å‰©ä½™æ—¶é—´
            time_remaining = None
            if link.expire_at:
                remaining = link.expire_at - datetime.now()
                if remaining.days > 0:
                    time_remaining = f"{remaining.days} å¤©"
                elif remaining.seconds > 3600:
                    time_remaining = f"{remaining.seconds // 3600} å°æ—¶"
                elif remaining.seconds > 60:
                    time_remaining = f"{remaining.seconds // 60} åˆ†é’Ÿ"
                else:
                    time_remaining = f"{remaining.seconds} ç§’"
            
            # è·å–ä¸»é¢˜å’ŒèƒŒæ™¯è®¾ç½®
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('share.html',
                                 token=token,
                                 file_name=file_name,
                                 file_type=file_type,
                                 file_size=file_size,
                                 created_at=created_at,
                                 expire_at=expire_at,
                                 time_remaining=time_remaining,
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color,
                                 is_batch=link.is_batch)
    except Exception as e:
        print(f"åˆ†äº«é¡µé¢åŠ è½½å¤±è´¥: {e}")
        import traceback
        traceback.print_exc()
        return "åŠ è½½å¤±è´¥", 500

# --- è·¯ç”±ï¼šåˆ†äº«æ–‡ä»¶ä¸‹è½½ ---
@app.route('/share-download/<token>')
def share_download(token):
    link = ShareLink.query.filter_by(token=token).first()
    
    if not link:
        return "é“¾æ¥ä¸å­˜åœ¨æˆ–å·²å¤±æ•ˆ", 404
        
    if link.expire_at and datetime.now() > link.expire_at:
        return "é“¾æ¥å·²è¿‡æœŸ", 403
        
    # ä¸‹è½½æ–‡ä»¶
    try:
        if link.is_batch:
            # æ‰¹é‡ä¸‹è½½ï¼šæ‰“åŒ…æˆ ZIP
            file_paths = link.file_path.split('|')
            
            # åˆ›å»ºä¸´æ—¶ ZIP æ–‡ä»¶
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            random_str = shortuuid.uuid()[:6]
            zip_name = f"batch_share_{timestamp}_{random_str}.zip"
            zip_path = os.path.join(FOLDERZIP_DIR, zip_name)
            
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for path in file_paths:
                    try:
                        abs_path = get_safe_path(path)
                        if not os.path.exists(abs_path):
                            continue
                        
                        if os.path.isdir(abs_path):
                            # æ·»åŠ æ–‡ä»¶å¤¹
                            folder_name = os.path.basename(abs_path)
                            for root, dirs, files in os.walk(abs_path):
                                for file in files:
                                    file_path = os.path.join(root, file)
                                    arcname = os.path.join(folder_name, os.path.relpath(file_path, abs_path))
                                    zipf.write(file_path, arcname)
                        else:
                            # æ·»åŠ æ–‡ä»¶
                            zipf.write(abs_path, os.path.basename(abs_path))
                    except Exception as e:
                        print(f"æ‰“åŒ…æ–‡ä»¶å¤±è´¥ {path}: {e}")
            
            return send_file(
                zip_path,
                mimetype='application/zip',
                as_attachment=True,
                download_name=f"æ‰¹é‡åˆ†äº«_{timestamp}.zip"
            )
        else:
            # å•ä¸ªæ–‡ä»¶ä¸‹è½½
            abs_path = get_safe_path(link.file_path)
            if os.path.isdir(abs_path):
                # å¦‚æœæ˜¯æ–‡ä»¶å¤¹ï¼Œæ‰“åŒ…ä¸º ZIP ä¸‹è½½
                folder_name = os.path.basename(abs_path)
                zip_path = zip_folder(abs_path, folder_name)
                
                return send_file(
                    zip_path,
                    mimetype='application/zip',
                    as_attachment=True,
                    download_name=f"{folder_name}.zip"
                )
            return send_file(abs_path, as_attachment=True)
    except Exception as e:
        print(f"åˆ†äº«ä¸‹è½½å¤±è´¥: {e}")
        import traceback
        traceback.print_exc()
        return "æ–‡ä»¶æºå·²è¢«åˆ é™¤æˆ–ä¸‹è½½å¤±è´¥", 404

# --- è·¯ç”±ï¼šæ‰¹é‡åˆ†äº«ä¸­çš„å•ä¸ªæ–‡ä»¶ä¸‹è½½ ---
@app.route('/share-download-single/<token>/<int:index>')
def share_download_single(token, index):
    link = ShareLink.query.filter_by(token=token).first()
    
    if not link:
        return "é“¾æ¥ä¸å­˜åœ¨æˆ–å·²å¤±æ•ˆ", 404
        
    if link.expire_at and datetime.now() > link.expire_at:
        return "é“¾æ¥å·²è¿‡æœŸ", 403
    
    if not link.is_batch:
        return "æ­¤é“¾æ¥ä¸æ˜¯æ‰¹é‡åˆ†äº«", 400
    
    try:
        file_paths = link.file_path.split('|')
        
        if index < 0 or index >= len(file_paths):
            return "æ–‡ä»¶ç´¢å¼•æ— æ•ˆ", 400
        
        path = file_paths[index]
        abs_path = get_safe_path(path)
        
        if not os.path.exists(abs_path):
            return "æ–‡ä»¶ä¸å­˜åœ¨æˆ–å·²è¢«åˆ é™¤", 404
        
        if os.path.isdir(abs_path):
            # å¦‚æœæ˜¯æ–‡ä»¶å¤¹ï¼Œæ‰“åŒ…ä¸º ZIP ä¸‹è½½
            folder_name = os.path.basename(abs_path)
            zip_path = zip_folder(abs_path, folder_name)
            
            return send_file(
                zip_path,
                mimetype='application/zip',
                as_attachment=True,
                download_name=f"{folder_name}.zip"
            )
        else:
            # ç›´æ¥ä¸‹è½½æ–‡ä»¶
            return send_file(abs_path, as_attachment=True)
    except Exception as e:
        print(f"å•ä¸ªæ–‡ä»¶ä¸‹è½½å¤±è´¥: {e}")
        import traceback
        traceback.print_exc()
        return "ä¸‹è½½å¤±è´¥", 500

# --- è·¯ç”±ï¼šåˆ†äº«æ–‡ä»¶é¢„è§ˆ ---
@app.route('/share-preview/<token>')
def share_preview(token):
    link = ShareLink.query.filter_by(token=token).first()
    
    if not link:
        return "é“¾æ¥ä¸å­˜åœ¨æˆ–å·²å¤±æ•ˆ", 404
        
    if link.expire_at and datetime.now() > link.expire_at:
        return "é“¾æ¥å·²è¿‡æœŸ", 403
    
    try:
        abs_path = get_safe_path(link.file_path)
        
        if not os.path.isfile(abs_path):
            return "æ–‡ä»¶ä¸å­˜åœ¨", 404
        
        filename = os.path.basename(abs_path)
        file_type = get_file_type(filename)
        
        if file_type not in ['image', 'video', 'audio', 'office', 'pdf']:
            return "æ­¤æ–‡ä»¶ç±»å‹ä¸æ”¯æŒé¢„è§ˆ", 400
        
        # Office æ–‡æ¡£å’Œ PDF ä½¿ç”¨æ–°çš„é¢„è§ˆé¡µé¢
        if file_type in ['office', 'pdf']:
            # ç”Ÿæˆå¯è®¿é—®çš„æ–‡ä»¶ URL
            file_url = request.host_url + 'share-file/' + token
            from urllib.parse import quote
            file_url = quote(file_url, safe=':/?&=')
            
            return render_template('document_preview.html', 
                                 file_path=link.file_path, 
                                 file_name=filename,
                                 file_type=file_type,
                                 file_url=file_url)
        
        # å›¾ç‰‡ã€è§†é¢‘å’ŒéŸ³é¢‘ä½¿ç”¨åŸæ¥çš„é¢„è§ˆé¡µé¢
        return render_template('preview.html', 
                             file_path=link.file_path, 
                             file_name=filename,
                             file_type=file_type)
    except Exception as e:
        print(f"é¢„è§ˆå¤±è´¥: {e}")
        return "é¢„è§ˆå¤±è´¥", 500

# --- è·¯ç”±ï¼šåˆ†äº«æ–‡ä»¶å†…å®¹ï¼ˆç”¨äºé¢„è§ˆï¼‰---
@app.route('/share-file/<token>')
def share_file(token):
    link = ShareLink.query.filter_by(token=token).first()
    
    if not link:
        abort(404)
        
    if link.expire_at and datetime.now() > link.expire_at:
        abort(403)
    
    try:
        abs_path = get_safe_path(link.file_path)
        
        if not os.path.isfile(abs_path):
            abort(404)
        
        filename = os.path.basename(abs_path)
        file_type = get_file_type(filename)
        
        if file_type == 'pdf':
            return send_file(abs_path, mimetype='application/pdf')
        elif file_type == 'office':
            ext = os.path.splitext(filename.lower())[1]
            mime_types = {
                '.doc': 'application/msword',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                '.xls': 'application/vnd.ms-excel',
                '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                '.ppt': 'application/vnd.ms-powerpoint',
                '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            }
            return send_file(abs_path, mimetype=mime_types.get(ext, 'application/octet-stream'))
        else:
            abort(404)
    except:
        abort(404)

# --- è·¯ç”±ï¼šåˆ†äº«æ–‡ä»¶ç¼©ç•¥å›¾ ---
@app.route('/share-thumbnail/<token>')
def share_thumbnail(token):
    link = ShareLink.query.filter_by(token=token).first()
    
    if not link:
        abort(404)
        
    if link.expire_at and datetime.now() > link.expire_at:
        abort(403)
    
    try:
        abs_path = get_safe_path(link.file_path)
        
        if not os.path.isfile(abs_path):
            abort(404)
        
        # åªä¸ºå›¾ç‰‡ç”Ÿæˆç¼©ç•¥å›¾
        if is_image(abs_path):
            try:
                img = Image.open(abs_path)
                # è½¬æ¢ RGBA åˆ° RGB
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                    img = background
                
                # ç”Ÿæˆç¼©ç•¥å›¾
                img.thumbnail((400, 400), Image.Resampling.LANCZOS)
                
                # ä¿å­˜åˆ°å†…å­˜
                img_io = io.BytesIO()
                img.save(img_io, 'JPEG', quality=85)
                img_io.seek(0)
                
                return send_file(img_io, mimetype='image/jpeg')
            except Exception as e:
                print(f"ç¼©ç•¥å›¾ç”Ÿæˆå¤±è´¥: {e}")
                abort(404)
        else:
            abort(404)
    except:
        abort(404)

# --- è·¯ç”±ï¼šç”Ÿæˆç¼©ç•¥å›¾ ---
@app.route('/thumbnail')
@login_required
def thumbnail():
    path = request.args.get('path')
    try:
        abs_path = get_safe_path(path)
        
        if not os.path.isfile(abs_path):
            abort(404)
        
        # åªä¸ºå›¾ç‰‡ç”Ÿæˆç¼©ç•¥å›¾
        if is_image(abs_path):
            try:
                img = Image.open(abs_path)
                # è½¬æ¢ RGBA åˆ° RGB
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                    img = background
                
                # ç”Ÿæˆç¼©ç•¥å›¾
                img.thumbnail((200, 200), Image.Resampling.LANCZOS)
                
                # ä¿å­˜åˆ°å†…å­˜
                img_io = io.BytesIO()
                img.save(img_io, 'JPEG', quality=85)
                img_io.seek(0)
                
                return send_file(img_io, mimetype='image/jpeg')
            except Exception as e:
                print(f"ç¼©ç•¥å›¾ç”Ÿæˆå¤±è´¥: {e}")
                abort(404)
        else:
            abort(404)
    except:
        abort(404)

# --- è·¯ç”±ï¼šæŸ¥çœ‹å‹ç¼©åŒ…å†…å®¹ ---
@app.route('/archive-view')
@login_required
def archive_view():
    path = request.args.get('path')
    try:
        abs_path = get_safe_path(path)
        if not os.path.isfile(abs_path):
            abort(404)
        
        if not is_archive(abs_path):
            abort(404)
        
        filename = os.path.basename(abs_path)
        
        # è·å–ä¸»é¢˜å’ŒèƒŒæ™¯è®¾ç½®
        theme = get_setting('theme', 'light')
        bg_type = get_setting('background_type', 'image')
        bg_image = get_setting('background_image', 'bg.png')
        bg_color = get_setting('background_color', '#667eea')
        
        return render_template('archive_view.html',
                             archive_path=path,
                             archive_name=filename,
                             theme=theme,
                             bg_type=bg_type,
                             bg_image=bg_image,
                             bg_color=bg_color)
    except:
        abort(404)

# --- è·¯ç”±ï¼šä»å‹ç¼©åŒ…ä¸‹è½½å•ä¸ªæ–‡ä»¶ ---
@app.route('/download-from-archive', methods=['POST'])
@login_required
def download_from_archive():
    archive_path = request.form.get('archive_path')
    file_name = request.form.get('file_name')
    
    try:
        abs_archive_path = get_safe_path(archive_path)
        
        if not os.path.exists(abs_archive_path):
            return "å‹ç¼©åŒ…ä¸å­˜åœ¨", 404
        
        file_ext = abs_archive_path.lower()
        
        # åˆ›å»ºä¸´æ—¶ç›®å½•
        import tempfile
        temp_dir = tempfile.mkdtemp()
        
        if file_ext.endswith('.zip'):
            import zipfile
            with zipfile.ZipFile(abs_archive_path, 'r') as zip_ref:
                # æŸ¥æ‰¾å¹¶æå–æ–‡ä»¶
                for file_info in zip_ref.filelist:
                    try:
                        filename = file_info.filename.encode('cp437').decode('utf-8')
                    except:
                        try:
                            filename = file_info.filename.encode('cp437').decode('gbk')
                        except:
                            filename = file_info.filename
                    
                    if filename == file_name:
                        extracted_path = zip_ref.extract(file_info, temp_dir)
                        
                        @after_this_request
                        def cleanup(response):
                            try:
                                shutil.rmtree(temp_dir)
                            except:
                                pass
                            return response
                        
                        return send_file(extracted_path, as_attachment=True, download_name=os.path.basename(file_name))
                        
        elif file_ext.endswith(('.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tbz2', '.tar.xz', '.txz')):
            import tarfile
            with tarfile.open(abs_archive_path, 'r:*') as tar_ref:
                member = tar_ref.getmember(file_name)
                tar_ref.extract(member, temp_dir)
                extracted_path = os.path.join(temp_dir, file_name)
                
                @after_this_request
                def cleanup(response):
                    try:
                        shutil.rmtree(temp_dir)
                    except:
                        pass
                    return response
                
                return send_file(extracted_path, as_attachment=True, download_name=os.path.basename(file_name))
        
        elif file_ext.endswith('.rar'):
            import rarfile
            with rarfile.RarFile(abs_archive_path, 'r') as rar_ref:
                # æå–æ–‡ä»¶
                rar_ref.extract(file_name, temp_dir)
                extracted_path = os.path.join(temp_dir, file_name)
                
                @after_this_request
                def cleanup(response):
                    try:
                        shutil.rmtree(temp_dir)
                    except:
                        pass
                    return response
                
                return send_file(extracted_path, as_attachment=True, download_name=os.path.basename(file_name))
                
        elif file_ext.endswith('.7z'):
            import py7zr
            with py7zr.SevenZipFile(abs_archive_path, 'r') as sz_ref:
                # æå–ç‰¹å®šæ–‡ä»¶
                sz_ref.extract(temp_dir, [file_name])
                extracted_path = os.path.join(temp_dir, file_name)
                
                @after_this_request
                def cleanup(response):
                    try:
                        shutil.rmtree(temp_dir)
                    except:
                        pass
                    return response
                
                return send_file(extracted_path, as_attachment=True, download_name=os.path.basename(file_name))
        
        return "æ–‡ä»¶æœªæ‰¾åˆ°", 404
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"ä¸‹è½½å¤±è´¥: {str(e)}", 500

# --- è¾…åŠ©å‡½æ•°ï¼šå°† Word æ–‡æ¡£è½¬æ¢ä¸º HTML ---
def convert_docx_to_html(docx_path):
    try:
        from docx import Document
        from docx.shared import Inches
        
        doc = Document(docx_path)
        html_content = ['<div class="document-content">']
        
        # å¤„ç†æ®µè½
        for para in doc.paragraphs:
            if para.text.strip():
                style = ''
                if para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    html_content.append(f'<h{level}>{para.text}</h{level}>')
                else:
                    # å¤„ç†æ®µè½æ ¼å¼
                    if para.alignment:
                        align_map = {0: 'left', 1: 'center', 2: 'right', 3: 'justify'}
                        style = f'text-align: {align_map.get(para.alignment, "left")};'
                    html_content.append(f'<p style="{style}">{para.text}</p>')
        
        # å¤„ç†è¡¨æ ¼
        for table in doc.tables:
            html_content.append('<table class="table table-bordered">')
            for row in table.rows:
                html_content.append('<tr>')
                for cell in row.cells:
                    html_content.append(f'<td>{cell.text}</td>')
                html_content.append('</tr>')
            html_content.append('</table>')
        
        html_content.append('</div>')
        return '\n'.join(html_content)
    except Exception as e:
        return f'<div class="alert alert-danger">Word æ–‡æ¡£è§£æå¤±è´¥: {str(e)}</div>'

# --- è¾…åŠ©å‡½æ•°ï¼šå°† Excel è¡¨æ ¼è½¬æ¢ä¸º HTML ---
def convert_xlsx_to_html(xlsx_path):
    try:
        from openpyxl import load_workbook
        
        wb = load_workbook(xlsx_path, data_only=True)
        html_content = ['<div class="document-content">']
        
        # å¤„ç†æ‰€æœ‰å·¥ä½œè¡¨
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            html_content.append(f'<h3>å·¥ä½œè¡¨: {sheet_name}</h3>')
            html_content.append('<div class="table-responsive">')
            html_content.append('<table class="table table-bordered table-striped">')
            
            # è·å–æœ‰æ•°æ®çš„åŒºåŸŸ
            max_row = sheet.max_row
            max_col = sheet.max_column
            
            # é™åˆ¶æ˜¾ç¤ºè¡Œæ•°ï¼Œé¿å…è¿‡å¤§
            display_rows = min(max_row, 1000)
            
            for row_idx, row in enumerate(sheet.iter_rows(max_row=display_rows, max_col=max_col), 1):
                html_content.append('<tr>')
                for cell in row:
                    value = cell.value if cell.value is not None else ''
                    # ç¬¬ä¸€è¡Œä½œä¸ºè¡¨å¤´
                    tag = 'th' if row_idx == 1 else 'td'
                    html_content.append(f'<{tag}>{value}</{tag}>')
                html_content.append('</tr>')
            
            if max_row > display_rows:
                html_content.append(f'<tr><td colspan="{max_col}" class="text-center text-muted">... è¿˜æœ‰ {max_row - display_rows} è¡Œæœªæ˜¾ç¤º ...</td></tr>')
            
            html_content.append('</table>')
            html_content.append('</div>')
        
        html_content.append('</div>')
        return '\n'.join(html_content)
    except Exception as e:
        return f'<div class="alert alert-danger">Excel è¡¨æ ¼è§£æå¤±è´¥: {str(e)}</div>'

# --- è¾…åŠ©å‡½æ•°ï¼šå°† PPT è½¬æ¢ä¸º HTML ---
def convert_pptx_to_html(pptx_path):
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        html_content = ['<div class="document-content presentation">']
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            html_content.append(f'<div class="slide" id="slide-{slide_idx}">')
            html_content.append(f'<h3>å¹»ç¯ç‰‡ {slide_idx}</h3>')
            html_content.append('<div class="slide-content">')
            
            # æå–æ–‡æœ¬å†…å®¹
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # åˆ¤æ–­æ˜¯å¦ä¸ºæ ‡é¢˜
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        html_content.append(f'<h4>{shape.text}</h4>')
                    else:
                        html_content.append(f'<p>{shape.text}</p>')
                
                # å¤„ç†è¡¨æ ¼
                if shape.has_table:
                    html_content.append('<table class="table table-bordered">')
                    for row in shape.table.rows:
                        html_content.append('<tr>')
                        for cell in row.cells:
                            html_content.append(f'<td>{cell.text}</td>')
                        html_content.append('</tr>')
                    html_content.append('</table>')
            
            html_content.append('</div>')
            html_content.append('</div>')
        
        html_content.append('</div>')
        return '\n'.join(html_content)
    except Exception as e:
        return f'<div class="alert alert-danger">PPT æ–‡æ¡£è§£æå¤±è´¥: {str(e)}</div>'

# --- è·¯ç”±ï¼šé¢„è§ˆæ–‡ä»¶ ---
@app.route('/preview')
@login_required
def preview():
    path = request.args.get('path')
    try:
        abs_path = get_safe_path(path)
        if not os.path.isfile(abs_path):
            abort(404)
        
        filename = os.path.basename(abs_path)
        file_type = get_file_type(filename)
        
        # Office æ–‡æ¡£è½¬æ¢ä¸º HTML é¢„è§ˆ
        if file_type == 'office':
            file_ext = os.path.splitext(filename.lower())[1]
            html_content = ''
            
            if file_ext in ['.docx', '.doc']:
                html_content = convert_docx_to_html(abs_path)
            elif file_ext in ['.xlsx', '.xls']:
                html_content = convert_xlsx_to_html(abs_path)
            elif file_ext in ['.pptx', '.ppt']:
                html_content = convert_pptx_to_html(abs_path)
            
            # è·å–ä¸»é¢˜å’ŒèƒŒæ™¯è®¾ç½®
            theme = get_setting('theme', 'light')
            bg_type = get_setting('background_type', 'image')
            bg_image = get_setting('background_image', 'bg.png')
            bg_color = get_setting('background_color', '#667eea')
            
            return render_template('office_preview.html',
                                 file_path=path,
                                 file_name=filename,
                                 file_type=file_type,
                                 html_content=html_content,
                                 theme=theme,
                                 bg_type=bg_type,
                                 bg_image=bg_image,
                                 bg_color=bg_color)
        
        # PDF ä½¿ç”¨æ–‡æ¡£é¢„è§ˆé¡µé¢
        elif file_type == 'pdf':
            return render_template('document_preview.html', 
                                 file_path=path, 
                                 file_name=filename,
                                 file_type=file_type,
                                 file_url='')
        
        # å›¾ç‰‡ã€è§†é¢‘å’ŒéŸ³é¢‘ä½¿ç”¨åŸæ¥çš„é¢„è§ˆé¡µé¢
        return render_template('preview.html', 
                             file_path=path, 
                             file_name=filename,
                             file_type=file_type)
    except:
        abort(404)

# --- è·¯ç”±ï¼šè·å–æ–‡ä»¶å†…å®¹ï¼ˆç”¨äºé¢„è§ˆï¼‰ ---
@app.route('/file')
@login_required
def get_file():
    path = request.args.get('path')
    try:
        abs_path = get_safe_path(path)
        if not os.path.isfile(abs_path):
            abort(404)
        
        # è·å–æ–‡ä»¶çš„ MIME ç±»å‹
        filename = os.path.basename(abs_path)
        file_type = get_file_type(filename)
        
        if file_type == 'image':
            return send_file(abs_path)
        elif file_type == 'video':
            return send_file(abs_path)
        elif file_type == 'audio':
            return send_file(abs_path)
        elif file_type == 'pdf':
            return send_file(abs_path, mimetype='application/pdf')
        elif file_type == 'office':
            # Office æ–‡æ¡£ç›´æ¥å‘é€
            ext = os.path.splitext(filename.lower())[1]
            mime_types = {
                '.doc': 'application/msword',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                '.xls': 'application/vnd.ms-excel',
                '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                '.ppt': 'application/vnd.ms-powerpoint',
                '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            }
            return send_file(abs_path, mimetype=mime_types.get(ext, 'application/octet-stream'))
        else:
            abort(404)
    except:
        abort(404)

# --- è·¯ç”±ï¼šä¸‹è½½æ™®é€šæ–‡ä»¶æˆ–æ–‡ä»¶å¤¹ ---
@app.route('/download')
@login_required
def download():
    path = request.args.get('path')
    try:
        abs_path = get_safe_path(path)
        if os.path.isdir(abs_path):
            # å¦‚æœæ˜¯æ–‡ä»¶å¤¹ï¼Œæ‰“åŒ…ä¸º ZIP ä¸‹è½½
            folder_name = os.path.basename(abs_path) or 'storage'
            zip_path = zip_folder(abs_path, folder_name)
            
            # ä¸å†ä¸‹è½½ååˆ é™¤ï¼Œä¿ç•™24å°æ—¶
            return send_file(
                zip_path,
                mimetype='application/zip',
                as_attachment=True,
                download_name=f"{folder_name}.zip"
            )
        return send_file(abs_path, as_attachment=True)
    except Exception as e:
        print(f"ä¸‹è½½å¤±è´¥: {e}")
        import traceback
        traceback.print_exc()
        return "æ–‡ä»¶ä¸å­˜åœ¨æˆ–ä¸‹è½½å¤±è´¥", 404

# --- æ¥å£ï¼šä¿®æ”¹å¯†ç  ---
@app.route('/api/change-password', methods=['POST'])
@login_required
def change_password():
    data = request.json
    old_password = data.get('old_password')
    new_password = data.get('new_password')
    
    if not old_password or not new_password:
        return jsonify({'status': 'error', 'msg': 'å¯†ç ä¸èƒ½ä¸ºç©º'})
    
    if not verify_password(old_password):
        return jsonify({'status': 'error', 'msg': 'åŸå¯†ç é”™è¯¯'})
    
    if len(new_password) < 6:
        return jsonify({'status': 'error', 'msg': 'æ–°å¯†ç è‡³å°‘6ä½'})
    
    # æ›´æ–°å¯†ç 
    new_hash = generate_password_hash(new_password)
    set_setting('password_hash', new_hash)
    
    return jsonify({'status': 'success', 'msg': 'å¯†ç ä¿®æ”¹æˆåŠŸ'})

# --- æ¥å£ï¼šåˆ‡æ¢ä¸»é¢˜ ---
@app.route('/api/toggle-theme', methods=['POST'])
@login_required
def toggle_theme():
    data = request.json
    theme = data.get('theme', 'light')
    
    if theme not in ['light', 'dark']:
        return jsonify({'status': 'error', 'msg': 'æ— æ•ˆçš„ä¸»é¢˜'})
    
    set_setting('theme', theme)
    return jsonify({'status': 'success', 'theme': theme})

# --- æ¥å£ï¼šè®¾ç½®å®‰å…¨é—®é¢˜ ---
@app.route('/api/set-security-question', methods=['POST'])
@login_required
def set_security_question():
    data = request.json
    question = data.get('question', '').strip()
    answer = data.get('answer', '').strip()
    
    if not question or not answer:
        return jsonify({'status': 'error', 'msg': 'é—®é¢˜å’Œç­”æ¡ˆä¸èƒ½ä¸ºç©º'})
    
    if len(answer) < 2:
        return jsonify({'status': 'error', 'msg': 'ç­”æ¡ˆè‡³å°‘éœ€è¦2ä¸ªå­—ç¬¦'})
    
    set_setting('security_question', question)
    set_setting('security_answer', answer)
    
    return jsonify({'status': 'success', 'msg': 'å®‰å…¨é—®é¢˜è®¾ç½®æˆåŠŸ'})

# --- æ¥å£ï¼šæ›´æ–°èƒŒæ™¯è®¾ç½® ---
@app.route('/api/update-background', methods=['POST'])
@login_required
def update_background():
    data = request.json
    bg_type = data.get('type', 'image')  # image æˆ– color
    bg_value = data.get('value', '')
    
    if bg_type == 'color':
        set_setting('background_type', 'color')
        set_setting('background_color', bg_value)
    elif bg_type == 'image':
        set_setting('background_type', 'image')
        set_setting('background_image', bg_value)
    else:
        return jsonify({'status': 'error', 'msg': 'æ— æ•ˆçš„èƒŒæ™¯ç±»å‹'})
    
    return jsonify({'status': 'success'})

# --- æ¥å£ï¼šä¸Šä¼ èƒŒæ™¯å›¾ç‰‡ ---
@app.route('/api/upload-background', methods=['POST'])
@login_required
def upload_background():
    file = request.files.get('file')
    
    if not file:
        return jsonify({'status': 'error', 'msg': 'æ²¡æœ‰æ–‡ä»¶'})
    
    # æ£€æŸ¥æ–‡ä»¶ç±»å‹
    if not file.filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp')):
        return jsonify({'status': 'error', 'msg': 'åªæ”¯æŒå›¾ç‰‡æ ¼å¼'})
    
    try:
        # ä¿å­˜æ–‡ä»¶
        filename = f"bg_{int(time.time())}.{file.filename.rsplit('.', 1)[1]}"
        filepath = os.path.join(STATIC_DIR, filename)
        file.save(filepath)
        
        # æ›´æ–°è®¾ç½®
        set_setting('background_type', 'image')
        set_setting('background_image', filename)
        
        return jsonify({'status': 'success', 'filename': filename})
    except Exception as e:
        return jsonify({'status': 'error', 'msg': str(e)})

# --- è·¯ç”±ï¼šè·å–å½“å‰è®¾ç½® ---
@app.route('/api/get-settings', methods=['GET'])
@login_required
def get_settings():
    return jsonify({
        'theme': get_setting('theme', 'light'),
        'background_type': get_setting('background_type', 'image'),
        'background_image': get_setting('background_image', 'bg.png'),
        'background_color': get_setting('background_color', '#667eea')
    })

# --- æ¥å£ï¼šæŸ¥çœ‹å‹ç¼©åŒ…å†…å®¹ ---
@app.route('/api/archive-content', methods=['POST'])
@login_required
def get_archive_content():
    data = request.json
    file_path = data.get('path')
    
    try:
        abs_file_path = get_safe_path(file_path)
        
        if not os.path.exists(abs_file_path):
            return jsonify({'status': 'error', 'msg': 'æ–‡ä»¶ä¸å­˜åœ¨'})
        
        if not os.path.isfile(abs_file_path):
            return jsonify({'status': 'error', 'msg': 'ä¸æ˜¯æœ‰æ•ˆçš„æ–‡ä»¶'})
        
        if not is_archive(abs_file_path):
            return jsonify({'status': 'error', 'msg': 'ä¸æ˜¯å‹ç¼©åŒ…æ–‡ä»¶'})
        
        file_list = []
        file_ext = abs_file_path.lower()
        
        if file_ext.endswith('.zip'):
            import zipfile
            with zipfile.ZipFile(abs_file_path, 'r') as zip_ref:
                for file_info in zip_ref.filelist:
                    try:
                        # å°è¯•è§£ç æ–‡ä»¶å
                        try:
                            filename = file_info.filename.encode('cp437').decode('utf-8')
                        except:
                            try:
                                filename = file_info.filename.encode('cp437').decode('gbk')
                            except:
                                filename = file_info.filename
                        
                        # è·³è¿‡ç©ºæ–‡ä»¶å¤¹ï¼ˆåªæœ‰è·¯å¾„åˆ†éš”ç¬¦çš„æ¡ç›®ï¼‰
                        if not filename or filename.strip('/') == '':
                            continue
                        
                        is_dir = filename.endswith('/')
                        size = file_info.file_size if not is_dir else 0
                        
                        # è·³è¿‡å¤§å°ä¸º0çš„ç›®å½•æ¡ç›®ï¼ˆç©ºæ–‡ä»¶å¤¹ï¼‰
                        if is_dir and size == 0:
                            continue
                        
                        # æ ¼å¼åŒ–å¤§å°
                        if size < 1024:
                            size_str = f"{size} B"
                        elif size < 1024 * 1024:
                            size_str = f"{size/1024:.2f} KB"
                        else:
                            size_str = f"{size/1024/1024:.2f} MB"
                        
                        file_list.append({
                            'name': filename.rstrip('/'),  # ç§»é™¤æœ«å°¾çš„æ–œæ 
                            'size': size_str,
                            'size_bytes': size,
                            'is_dir': is_dir,
                            'compressed_size': file_info.compress_size
                        })
                    except Exception as e:
                        print(f"å¤„ç†æ–‡ä»¶ä¿¡æ¯å¤±è´¥: {e}")
                        continue
                        
        elif file_ext.endswith(('.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tbz2', '.tar.xz', '.txz')):
            import tarfile
            with tarfile.open(abs_file_path, 'r:*') as tar_ref:
                for member in tar_ref.getmembers():
                    # è·³è¿‡ç©ºæ–‡ä»¶å¤¹
                    if not member.name or member.name.strip('/') == '':
                        continue
                    
                    is_dir = member.isdir()
                    size = member.size if not is_dir else 0
                    
                    # è·³è¿‡å¤§å°ä¸º0çš„ç›®å½•æ¡ç›®ï¼ˆç©ºæ–‡ä»¶å¤¹ï¼‰
                    if is_dir and size == 0:
                        continue
                    
                    # æ ¼å¼åŒ–å¤§å°
                    if size < 1024:
                        size_str = f"{size} B"
                    elif size < 1024 * 1024:
                        size_str = f"{size/1024:.2f} KB"
                    else:
                        size_str = f"{size/1024/1024:.2f} MB"
                    
                    file_list.append({
                        'name': member.name.rstrip('/'),
                        'size': size_str,
                        'size_bytes': size,
                        'is_dir': is_dir,
                        'compressed_size': 0
                    })
                    
        elif file_ext.endswith('.rar'):
            try:
                import rarfile
                with rarfile.RarFile(abs_file_path, 'r') as rar_ref:
                    for file_info in rar_ref.infolist():
                        # è·³è¿‡ç©ºæ–‡ä»¶å¤¹
                        if not file_info.filename or file_info.filename.strip('/') == '':
                            continue
                        
                        is_dir = file_info.isdir()
                        size = file_info.file_size if not is_dir else 0
                        
                        # è·³è¿‡å¤§å°ä¸º0çš„ç›®å½•æ¡ç›®ï¼ˆç©ºæ–‡ä»¶å¤¹ï¼‰
                        if is_dir and size == 0:
                            continue
                        
                        # æ ¼å¼åŒ–å¤§å°
                        if size < 1024:
                            size_str = f"{size} B"
                        elif size < 1024 * 1024:
                            size_str = f"{size/1024:.2f} KB"
                        else:
                            size_str = f"{size/1024/1024:.2f} MB"
                        
                        file_list.append({
                            'name': file_info.filename.rstrip('/'),
                            'size': size_str,
                            'size_bytes': size,
                            'is_dir': is_dir,
                            'compressed_size': file_info.compress_size
                        })
            except ImportError:
                return jsonify({'status': 'error', 'msg': 'RAR æ ¼å¼éœ€è¦å®‰è£… rarfile åº“ï¼Œè¯·è¿è¡Œ: pip install rarfile'})
            except Exception as e:
                return jsonify({'status': 'error', 'msg': f'RAR æ–‡ä»¶è¯»å–å¤±è´¥: {str(e)}'})
                
        elif file_ext.endswith('.7z'):
            try:
                import py7zr
                with py7zr.SevenZipFile(abs_file_path, 'r') as sz_ref:
                    for name, file_info in sz_ref.list():
                        # è·³è¿‡ç©ºæ–‡ä»¶å¤¹
                        if not name or name.strip('/') == '':
                            continue
                        
                        is_dir = file_info.is_directory
                        size = file_info.uncompressed if not is_dir else 0
                        
                        # è·³è¿‡å¤§å°ä¸º0çš„ç›®å½•æ¡ç›®ï¼ˆç©ºæ–‡ä»¶å¤¹ï¼‰
                        if is_dir and size == 0:
                            continue
                        
                        # æ ¼å¼åŒ–å¤§å°
                        if size < 1024:
                            size_str = f"{size} B"
                        elif size < 1024 * 1024:
                            size_str = f"{size/1024:.2f} KB"
                        else:
                            size_str = f"{size/1024/1024:.2f} MB"
                        
                        file_list.append({
                            'name': name.rstrip('/'),
                            'size': size_str,
                            'size_bytes': size,
                            'is_dir': is_dir,
                            'compressed_size': file_info.compressed if hasattr(file_info, 'compressed') else 0
                        })
            except ImportError:
                return jsonify({'status': 'error', 'msg': '7Z æ ¼å¼éœ€è¦å®‰è£… py7zr åº“ï¼Œè¯·è¿è¡Œ: pip install py7zr'})
            except Exception as e:
                return jsonify({'status': 'error', 'msg': f'7Z æ–‡ä»¶è¯»å–å¤±è´¥: {str(e)}'})
                
        else:
            return jsonify({'status': 'error', 'msg': 'ä¸æ”¯æŒæŸ¥çœ‹æ­¤å‹ç¼©æ ¼å¼çš„å†…å®¹'})
        
        # è¿‡æ»¤æ‰é‡å¤çš„æ¡ç›®ï¼ˆæœ‰äº›å‹ç¼©åŒ…ä¼šåŒæ—¶åŒ…å«æ–‡ä»¶å’Œå…¶çˆ¶ç›®å½•ï¼‰
        seen_names = set()
        filtered_list = []
        for item in file_list:
            if item['name'] not in seen_names:
                seen_names.add(item['name'])
                filtered_list.append(item)
        
        file_list = filtered_list
        
        # è®¡ç®—æ€»å¤§å°å’Œæ–‡ä»¶æ•°é‡
        total_size = sum(f['size_bytes'] for f in file_list if not f['is_dir'])
        file_count = len([f for f in file_list if not f['is_dir']])
        folder_count = len([f for f in file_list if f['is_dir']])
        
        if total_size < 1024:
            total_size_str = f"{total_size} B"
        elif total_size < 1024 * 1024:
            total_size_str = f"{total_size/1024:.2f} KB"
        else:
            total_size_str = f"{total_size/1024/1024:.2f} MB"
        
        return jsonify({
            'status': 'success',
            'files': file_list,
            'total_size': total_size_str,
            'file_count': file_count,
            'folder_count': folder_count
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'msg': f'è¯»å–å¤±è´¥: {str(e)}'})

# --- æ¥å£ï¼šä»å‹ç¼©åŒ…ä¸­æå–å•ä¸ªæ–‡ä»¶ ---
@app.route('/api/extract-file', methods=['POST'])
@login_required
def extract_single_file():
    data = request.json
    archive_path = data.get('archive_path')
    file_name = data.get('file_name')
    
    try:
        abs_archive_path = get_safe_path(archive_path)
        
        if not os.path.exists(abs_archive_path):
            return jsonify({'status': 'error', 'msg': 'å‹ç¼©åŒ…ä¸å­˜åœ¨'})
        
        file_ext = abs_archive_path.lower()
        
        # åˆ›å»ºä¸´æ—¶æ–‡ä»¶
        import tempfile
        temp_dir = tempfile.mkdtemp()
        
        if file_ext.endswith('.zip'):
            import zipfile
            with zipfile.ZipFile(abs_archive_path, 'r') as zip_ref:
                # æŸ¥æ‰¾æ–‡ä»¶
                for file_info in zip_ref.filelist:
                    try:
                        filename = file_info.filename.encode('cp437').decode('utf-8')
                    except:
                        try:
                            filename = file_info.filename.encode('cp437').decode('gbk')
                        except:
                            filename = file_info.filename
                    
                    if filename == file_name:
                        # æå–æ–‡ä»¶åˆ°ä¸´æ—¶ç›®å½•
                        extracted_path = zip_ref.extract(file_info, temp_dir)
                        
                        # è¿”å›æ–‡ä»¶è·¯å¾„ä¾›ä¸‹è½½
                        return jsonify({
                            'status': 'success',
                            'temp_path': extracted_path,
                            'file_name': os.path.basename(file_name)
                        })
                        
        elif file_ext.endswith(('.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tbz2', '.tar.xz', '.txz')):
            import tarfile
            with tarfile.open(abs_archive_path, 'r:*') as tar_ref:
                member = tar_ref.getmember(file_name)
                tar_ref.extract(member, temp_dir)
                extracted_path = os.path.join(temp_dir, file_name)
                
                return jsonify({
                    'status': 'success',
                    'temp_path': extracted_path,
                    'file_name': os.path.basename(file_name)
                })
        
        return jsonify({'status': 'error', 'msg': 'æ–‡ä»¶æœªæ‰¾åˆ°'})
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'msg': f'æå–å¤±è´¥: {str(e)}'})

# --- æ¥å£ï¼šè§£å‹å‹ç¼©åŒ… ---
@app.route('/api/extract', methods=['POST'])
@login_required
def extract_archive():
    data = request.json
    file_path = data.get('path')  # å‹ç¼©åŒ…è·¯å¾„
    extract_to = data.get('extract_to', '')  # è§£å‹åˆ°çš„ç›®æ ‡æ–‡ä»¶å¤¹ï¼Œé»˜è®¤ä¸ºå½“å‰ç›®å½•
    
    try:
        abs_file_path = get_safe_path(file_path)
        
        # æ£€æŸ¥æ–‡ä»¶æ˜¯å¦å­˜åœ¨
        if not os.path.exists(abs_file_path):
            return jsonify({'status': 'error', 'msg': 'æ–‡ä»¶ä¸å­˜åœ¨'})
        
        if not os.path.isfile(abs_file_path):
            return jsonify({'status': 'error', 'msg': 'ä¸æ˜¯æœ‰æ•ˆçš„æ–‡ä»¶'})
        
        # æ£€æŸ¥æ˜¯å¦ä¸ºå‹ç¼©åŒ…
        if not is_archive(abs_file_path):
            return jsonify({'status': 'error', 'msg': 'ä¸æ˜¯æ”¯æŒçš„å‹ç¼©åŒ…æ ¼å¼'})
        
        # ç¡®å®šè§£å‹ç›®æ ‡ç›®å½•
        if extract_to:
            abs_extract_dir = get_safe_path(extract_to)
        else:
            # é»˜è®¤è§£å‹åˆ°å‹ç¼©åŒ…æ‰€åœ¨ç›®å½•
            abs_extract_dir = os.path.dirname(abs_file_path)
        
        # åˆ›å»ºä»¥å‹ç¼©åŒ…åç§°å‘½åçš„æ–‡ä»¶å¤¹
        archive_name = os.path.splitext(os.path.basename(abs_file_path))[0]
        # å¤„ç† .tar.gz ç­‰åŒæ‰©å±•å
        if archive_name.endswith('.tar'):
            archive_name = os.path.splitext(archive_name)[0]
        
        extract_folder = os.path.join(abs_extract_dir, archive_name)
        
        # å¦‚æœç›®æ ‡æ–‡ä»¶å¤¹å·²å­˜åœ¨ï¼Œæ·»åŠ æ•°å­—åç¼€
        counter = 1
        original_extract_folder = extract_folder
        while os.path.exists(extract_folder):
            extract_folder = f"{original_extract_folder}_{counter}"
            counter += 1
        
        os.makedirs(extract_folder, exist_ok=True)
        
        # æ ¹æ®æ–‡ä»¶ç±»å‹è§£å‹
        file_ext = abs_file_path.lower()
        
        if file_ext.endswith('.zip'):
            # è§£å‹ ZIP æ–‡ä»¶
            import zipfile
            with zipfile.ZipFile(abs_file_path, 'r') as zip_ref:
                # å¤„ç†ä¸­æ–‡æ–‡ä»¶åç¼–ç é—®é¢˜
                for file_info in zip_ref.filelist:
                    try:
                        # å°è¯•ä½¿ç”¨ UTF-8 è§£ç 
                        file_info.filename = file_info.filename.encode('cp437').decode('utf-8')
                    except:
                        try:
                            # å°è¯•ä½¿ç”¨ GBK è§£ç 
                            file_info.filename = file_info.filename.encode('cp437').decode('gbk')
                        except:
                            pass
                zip_ref.extractall(extract_folder)
                
        elif file_ext.endswith(('.tar', '.tar.gz', '.tgz', '.tar.bz2', '.tbz2', '.tar.xz', '.txz')):
            # è§£å‹ TAR æ–‡ä»¶
            import tarfile
            with tarfile.open(abs_file_path, 'r:*') as tar_ref:
                tar_ref.extractall(extract_folder)
                
        elif file_ext.endswith('.gz') and not file_ext.endswith('.tar.gz'):
            # è§£å‹å•ä¸ª GZ æ–‡ä»¶
            import gzip
            output_file = os.path.join(extract_folder, os.path.splitext(os.path.basename(abs_file_path))[0])
            with gzip.open(abs_file_path, 'rb') as f_in:
                with open(output_file, 'wb') as f_out:
                    shutil.copyfileobj(f_in, f_out)
                    
        elif file_ext.endswith('.bz2') and not file_ext.endswith('.tar.bz2'):
            # è§£å‹å•ä¸ª BZ2 æ–‡ä»¶
            import bz2
            output_file = os.path.join(extract_folder, os.path.splitext(os.path.basename(abs_file_path))[0])
            with bz2.open(abs_file_path, 'rb') as f_in:
                with open(output_file, 'wb') as f_out:
                    shutil.copyfileobj(f_in, f_out)
                    
        elif file_ext.endswith('.rar'):
            # è§£å‹ RAR æ–‡ä»¶
            try:
                import rarfile
                with rarfile.RarFile(abs_file_path, 'r') as rar_ref:
                    rar_ref.extractall(extract_folder)
            except ImportError:
                return jsonify({'status': 'error', 'msg': 'RAR æ ¼å¼éœ€è¦å®‰è£… rarfile åº“ï¼Œè¯·è¿è¡Œ: pip install rarfile'})
            except Exception as e:
                return jsonify({'status': 'error', 'msg': f'RAR è§£å‹å¤±è´¥: {str(e)}'})
            
        elif file_ext.endswith('.7z'):
            # è§£å‹ 7Z æ–‡ä»¶
            try:
                import py7zr
                with py7zr.SevenZipFile(abs_file_path, 'r') as sz_ref:
                    sz_ref.extractall(extract_folder)
            except ImportError:
                return jsonify({'status': 'error', 'msg': '7Z æ ¼å¼éœ€è¦å®‰è£… py7zr åº“ï¼Œè¯·è¿è¡Œ: pip install py7zr'})
            except Exception as e:
                return jsonify({'status': 'error', 'msg': f'7Z è§£å‹å¤±è´¥: {str(e)}'})
            
        else:
            return jsonify({'status': 'error', 'msg': 'ä¸æ”¯æŒçš„å‹ç¼©æ ¼å¼'})
        
        return jsonify({
            'status': 'success', 
            'msg': f'è§£å‹æˆåŠŸï¼Œæ–‡ä»¶å·²è§£å‹åˆ°: {os.path.basename(extract_folder)}'
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'msg': f'è§£å‹å¤±è´¥: {str(e)}'})

# --- æ¥å£ï¼šæ¸…ç©ºç¼“å­˜ ---
@app.route('/api/clear-cache', methods=['POST'])
@login_required
def clear_cache():
    try:
        # æ¸…ç©º folderzip ç›®å½•ä¸­çš„æ‰€æœ‰ ZIP æ–‡ä»¶
        deleted_count = 0
        if os.path.exists(FOLDERZIP_DIR):
            for filename in os.listdir(FOLDERZIP_DIR):
                filepath = os.path.join(FOLDERZIP_DIR, filename)
                if os.path.isfile(filepath) and filename.endswith('.zip'):
                    os.remove(filepath)
                    deleted_count += 1
        
        # æ¸…ç©ºæ—§çš„èƒŒæ™¯å›¾ç‰‡ï¼ˆä¿ç•™å½“å‰ä½¿ç”¨çš„ï¼‰
        current_bg = get_setting('background_image', 'bg.png')
        if os.path.exists(STATIC_DIR):
            for filename in os.listdir(STATIC_DIR):
                if filename.startswith('bg_') and filename != current_bg:
                    filepath = os.path.join(STATIC_DIR, filename)
                    if os.path.isfile(filepath):
                        os.remove(filepath)
                        deleted_count += 1
        
        return jsonify({
            'status': 'success', 
            'msg': f'å·²æ¸…ç©ºç¼“å­˜ï¼Œåˆ é™¤äº† {deleted_count} ä¸ªä¸´æ—¶æ–‡ä»¶'
        })
    except Exception as e:
        return jsonify({'status': 'error', 'msg': str(e)})

# --- æ¥å£ï¼šæ¸…ç©ºæ‰€æœ‰æ•°æ® ---
@app.route('/api/clear-all-data', methods=['POST'])
@login_required
def clear_all_data():
    data = request.json
    confirm_text = data.get('confirm', '')
    
    # éœ€è¦è¾“å…¥ç¡®è®¤æ–‡æœ¬
    if confirm_text != 'DELETE ALL':
        return jsonify({'status': 'error', 'msg': 'è¯·è¾“å…¥æ­£ç¡®çš„ç¡®è®¤æ–‡æœ¬'})
    
    try:
        # 1. åˆ é™¤æ‰€æœ‰äº‘ç›˜æ–‡ä»¶
        if os.path.exists(STORAGE_DIR):
            shutil.rmtree(STORAGE_DIR)
            os.makedirs(STORAGE_DIR)
        
        # 2. åˆ é™¤æ‰€æœ‰ ZIP æ–‡ä»¶
        if os.path.exists(FOLDERZIP_DIR):
            shutil.rmtree(FOLDERZIP_DIR)
            os.makedirs(FOLDERZIP_DIR)
        
        # 3. åˆ é™¤æ‰€æœ‰ä¸Šä¼ çš„èƒŒæ™¯å›¾ç‰‡
        if os.path.exists(STATIC_DIR):
            for filename in os.listdir(STATIC_DIR):
                if filename.startswith('bg_'):
                    filepath = os.path.join(STATIC_DIR, filename)
                    if os.path.isfile(filepath):
                        os.remove(filepath)
        
        # 4. æ¸…ç©ºæ•°æ®åº“ä¸­çš„åˆ†äº«é“¾æ¥
        ShareLink.query.delete()
        
        # 5. é‡ç½®æ‰€æœ‰è®¾ç½®ä¸ºé»˜è®¤å€¼
        default_hash = generate_password_hash(DEFAULT_PASSWORD)
        set_setting('password_hash', default_hash)
        set_setting('theme', 'light')
        set_setting('background_type', 'image')
        set_setting('background_image', 'bg.png')
        set_setting('background_color', '#667eea')
        
        db.session.commit()
        
        # 6. æ¸…é™¤å½“å‰ sessionï¼Œå¼ºåˆ¶é‡æ–°ç™»å½•
        session.clear()
        
        return jsonify({
            'status': 'success', 
            'msg': 'æ‰€æœ‰æ•°æ®å·²æ¸…ç©ºï¼Œå¯†ç å·²é‡ç½®ä¸ºé»˜è®¤å¯†ç '
        })
    except Exception as e:
        db.session.rollback()
        return jsonify({'status': 'error', 'msg': str(e)})

# --- ä¸»ç¨‹åºå…¥å£

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000)
